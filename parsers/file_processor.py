import os
import base64
import io
from typing import Any


class FileProcessor:
    """Verwerkt PDF, PPTX, XLSX, DOCX en afbeeldingen naar gestructureerde tekst."""

    MAX_TEXT_PER_FILE = 4000   # tekens per bestand
    MAX_ROWS_PER_SHEET = 60    # rijen per Excel-tabblad
    MAX_PDF_PAGES = 15
    MAX_IMAGES = 4             # max afbeeldingen naar Claude

    def process_single_file(self, filepath: str) -> dict[str, Any] | None:
        """Process one file and return a single result dict, or None on failure."""
        results = self.process_files([filepath])
        return results[0] if results else None

    def process_files(self, filepaths: list[str]) -> list[dict[str, Any]]:
        results = []
        image_count = 0

        for filepath in filepaths:
            if not os.path.exists(filepath):
                continue

            ext = os.path.splitext(filepath)[1].lower()
            filename = os.path.basename(filepath)

            try:
                if ext == '.pdf':
                    data = self._process_pdf(filepath, filename)
                elif ext == '.pptx':
                    data = self._process_pptx(filepath, filename)
                elif ext in ('.xlsx', '.xls'):
                    data = self._process_xlsx(filepath, filename)
                elif ext == '.docx':
                    data = self._process_docx(filepath, filename)
                elif ext in ('.png', '.jpg', '.jpeg', '.gif', '.webp'):
                    if image_count < self.MAX_IMAGES:
                        data = self._process_image(filepath, filename)
                        image_count += 1
                    else:
                        data = {
                            'filename': filename, 'type': 'image',
                            'text': f'[{filename}: overgeslagen — max {self.MAX_IMAGES} afbeeldingen bereikt]',
                            'image_data': None, 'metadata': {}
                        }
                else:
                    continue

                results.append(data)

            except Exception as exc:
                results.append({
                    'filename': filename,
                    'type': ext.lstrip('.'),
                    'text': f'[Verwerkingsfout: {exc}]',
                    'image_data': None,
                    'metadata': {},
                    'error': str(exc),
                })

        return results

    # ------------------------------------------------------------------ #
    #  Helpers                                                             #
    # ------------------------------------------------------------------ #

    def _truncate(self, text: str) -> str:
        if len(text) > self.MAX_TEXT_PER_FILE:
            return (text[:self.MAX_TEXT_PER_FILE]
                    + f'\n\n[... Tekst ingekort: {len(text)} → {self.MAX_TEXT_PER_FILE} tekens ...]')
        return text

    # ------------------------------------------------------------------ #
    #  File-type parsers                                                   #
    # ------------------------------------------------------------------ #

    def _process_pdf(self, filepath: str, filename: str) -> dict:
        import pdfplumber

        parts = []
        page_count = 0

        with pdfplumber.open(filepath) as pdf:
            page_count = len(pdf.pages)
            for i, page in enumerate(pdf.pages[:self.MAX_PDF_PAGES]):
                text = page.extract_text() or ''
                if text.strip():
                    parts.append(f'[Pagina {i + 1}]\n{text}')

        full_text = '\n\n'.join(parts)
        is_scan = len(full_text.strip()) < 200 and page_count > 0

        return {
            'filename': filename,
            'type': 'pdf',
            'text': self._truncate(full_text) if full_text.strip()
                    else '[PDF zonder extraheerbare tekst — waarschijnlijk een scan of beveiligd bestand]',
            'image_data': None,
            'is_scan': is_scan,
            'metadata': {'pages': page_count},
        }

    def _process_pptx(self, filepath: str, filename: str) -> dict:
        from pptx import Presentation

        prs = Presentation(filepath)
        slides_content = []

        for i, slide in enumerate(prs.slides):
            slide_parts = []

            # Titel
            if slide.shapes.title and slide.shapes.title.text.strip():
                slide_parts.append(f'🔹 {slide.shapes.title.text.strip()}')

            # Alle shapes
            for shape in slide.shapes:
                if shape is slide.shapes.title:
                    continue
                if hasattr(shape, 'text') and shape.text.strip():
                    slide_parts.append(shape.text.strip())

            # Notities (speaker notes)
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    slide_parts.append(f'[Speaker notes]: {notes}')

            if slide_parts:
                slides_content.append(f'[Slide {i + 1}]\n' + '\n'.join(slide_parts))

        return {
            'filename': filename,
            'type': 'pptx',
            'text': self._truncate('\n\n'.join(slides_content)),
            'image_data': None,
            'metadata': {'slides': len(prs.slides)},
        }

    def _process_xlsx(self, filepath: str, filename: str) -> dict:
        import openpyxl

        wb = openpyxl.load_workbook(filepath, data_only=True, read_only=True)
        all_sheets = []

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows_text = [f'[Tabblad: {sheet_name}]']
            row_count = 0

            for row in ws.iter_rows(values_only=True):
                if row_count >= self.MAX_ROWS_PER_SHEET:
                    rows_text.append(f'[... nog meer rijen ...]')
                    break
                if any(v is not None for v in row):
                    rows_text.append(' | '.join('' if v is None else str(v) for v in row))
                    row_count += 1

            if row_count:
                all_sheets.append('\n'.join(rows_text))

        wb.close()

        return {
            'filename': filename,
            'type': 'xlsx',
            'text': self._truncate('\n\n'.join(all_sheets)),
            'image_data': None,
            'metadata': {'sheets': list(wb.sheetnames) if hasattr(wb, 'sheetnames') else []},
        }

    def _process_docx(self, filepath: str, filename: str) -> dict:
        from docx import Document

        doc = Document(filepath)
        parts = []

        for para in doc.paragraphs:
            if not para.text.strip():
                continue
            if para.style.name.startswith('Heading'):
                level_str = para.style.name.replace('Heading', '').strip()
                level = int(level_str) if level_str.isdigit() else 2
                parts.append('\n' + '#' * level + ' ' + para.text)
            else:
                parts.append(para.text)

        for i, table in enumerate(doc.tables):
            parts.append(f'\n[Tabel {i + 1}]')
            for row in table.rows:
                row_text = ' | '.join(c.text.strip() for c in row.cells)
                if row_text.replace('|', '').strip():
                    parts.append(row_text)

        return {
            'filename': filename,
            'type': 'docx',
            'text': self._truncate('\n'.join(parts)),
            'image_data': None,
            'metadata': {},
        }

    def _process_image(self, filepath: str, filename: str) -> dict:
        from PIL import Image

        with Image.open(filepath) as img:
            if img.mode not in ('RGB', 'RGBA'):
                img = img.convert('RGB')
            elif img.mode == 'RGBA':
                bg = Image.new('RGB', img.size, (255, 255, 255))
                bg.paste(img, mask=img.split()[3])
                img = bg

            # Verklein indien te groot
            max_dim = 1600
            if max(img.size) > max_dim:
                ratio = max_dim / max(img.size)
                img = img.resize(
                    (int(img.size[0] * ratio), int(img.size[1] * ratio)),
                    Image.LANCZOS
                )

            buf = io.BytesIO()
            img.save(buf, format='JPEG', quality=85)
            b64 = base64.standard_b64encode(buf.getvalue()).decode('utf-8')

        return {
            'filename': filename,
            'type': 'image',
            'text': f'[Afbeelding: {filename}]',
            'image_data': {'base64': b64, 'media_type': 'image/jpeg'},
            'metadata': {},
        }
